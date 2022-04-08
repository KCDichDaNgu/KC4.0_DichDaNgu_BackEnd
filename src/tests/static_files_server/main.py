import unittest2
import json
import requests 
import os
import docx 
import pptx
import json
import tqdm
import random

from openpyxl import load_workbook


def url_gen(path):
    form = "http://localhost:8001/{}".format(path)
    return form

def read_and_save(URL):
    the_book = requests.get(URL, stream=True)

    if the_book.status_code != 200:
        try:
            the_book.raise_for_status() 
        except requests.HTTPError as exception:
            return None, the_book.status_code, exception

    book_name = 'src/tests/static_files_server/static' + '/' + URL.split('/')[-1]

    with open(book_name, 'wb') as f:
        for chunk in the_book.iter_content(1024 * 1024 * 2):  # 2 MB chunks
            f.write(chunk)

    return book_name, the_book.status_code, 'OK'


def get_all_text_xlsx(filename):
    wb = load_workbook(filename, data_only=True)
    fullText = []

    for sheetname in wb.sheetnames:
        ws = wb[sheetname]
        all_rows = list(ws.rows)

        for row in all_rows:
            for cell in row:
                fullText.append(cell.value)
    return '\n'.join(fullText)

def get_all_text_docx(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)


def get_all_text_pptx(filename):
    prs = pptx.Presentation(filename)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text) 
    return '\n'.join(fullText)


def get_all_text_json(filename):
    f = open(filename)

    data = json.load(f)
    data = json.dumps(data)

    f.close()
    return data


def mutating_path(path):
    components = path.split('/')
    components[-1]

    if random.random() < 0.5:
        pos = random.randint(1, len(components[-2]) - 1)
        if components[-2][pos] == 'a':
            components[-2] = components[-2][:pos] + 'b' + components[-2][pos + 1:]
        else:
            components[-2] = components[-2][:pos] + 'a' + components[-2][pos + 1:] 

        pos = random.randint(1, len(components[-1]) - 1)
        if components[-1][pos] == 'a':
            components[-1] = components[-2][:pos] + 'b' + components[-1][pos + 1:]
        else:
            components[-1] = components[-2][:pos] + 'a' + components[-1][pos + 1:] 
    else:
        if random.random() < 0.5:
            pos = random.randint(1, len(components[-2]) - 1)
            if components[-2][pos] == 'a':
                components[-2] = components[-2][:pos] + 'b' + components[-2][pos + 1:]
            else:
                components[-2] = components[-2][:pos] + 'a' + components[-2][pos + 1:] 
        else:
            pos = random.randint(1, len(components[-1]) - 1)
            if components[-1][pos] == 'a':
                components[-1] = components[-2][:pos] + 'b' + components[-1][pos + 1:]
            else:
                components[-1] = components[-2][:pos] + 'a' + components[-1][pos + 1:] 
    return '/'.join(components)


class TestStaticFilesServer(unittest2.TestCase):
    def __init__(self, *args, **kwargs):
        super(TestStaticFilesServer, self).__init__(*args, **kwargs)
        
        self.path = "static"
        self.baseURL = "http://localhost:8001/static"

        self.filepaths = []

        for task_type in os.listdir(self.path):
            path1 = os.path.join(self.path, task_type)

            for requestId in os.listdir(path1):
                path2 = os.path.join(path1, requestId)

                if not os.path.isdir(path2):    
                    self.filepaths.append(path2)
                else:
                    for filename in os.listdir(path2):
                        path3 = os.path.join(path2, filename)

                        if path3.split('.')[-1] not in ['json', 'pptx', 'docx', 'xlsx']:
                            continue

                        self.filepaths.append(path3)
                        
        print("Total number of files for testing: ", len(self.filepaths))

    def test_ok_case(self):
        for _, path in tqdm.tqdm(enumerate(self.filepaths)):
            saved_path, status_code, message = read_and_save(url_gen(path))
            
            self.assertEqual(status_code, 200, "failed at {}".format(path))

            if path.split('.')[-1] == 'json':
                self.assertEqual(get_all_text_json(saved_path), get_all_text_json(path), "failed at {}".format(path))
            elif path.split('.')[-1] == 'docx':
                self.assertEqual(get_all_text_docx(saved_path), get_all_text_docx(path), "failed at {}".format(path))
            elif path.split('.')[-1] == 'pptx':
                self.assertEqual(get_all_text_pptx(saved_path), get_all_text_pptx(path), "failed at {}".format(path))
            elif path.split('.')[-1] == 'xlsx':
                self.assertEqual(get_all_text_xlsx(saved_path), get_all_text_xlsx(path), "failed at {}".format(path))

            os.remove(saved_path)


    def test_not_ok_test(self):
        mutated_filepaths = []
        for filepaths in self.filepaths:
            mutated_filepaths.append(mutating_path(filepaths))

        for _, path in tqdm.tqdm(enumerate(mutated_filepaths)):
            saved_path, status_code, message = read_and_save(url_gen(path))
            
            self.assertEqual(status_code, 404, "failed at {}".format(path))
            self.assertEqual(str(message), "404 Client Error: Not Found for url: {}".format(url_gen(path)), "failed at {}".format(path))
                


if __name__ == "__main__":
    unittest2.main(verbosity=2)
